from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Pt


@dataclass
class SlideSection:
    heading: str
    title: str
    body_lines: list[str]


def _clean_line(line: str) -> str:
    text = line.strip()
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = text.replace("`", "")
    return text.strip()


def parse_markdown_sections(markdown_text: str) -> list[SlideSection]:
    lines = markdown_text.splitlines()
    sections: list[tuple[str, list[str]]] = []
    current_heading: str | None = None
    current_block: list[str] = []

    for raw in lines:
        line = raw.rstrip("\n")
        if line.startswith("## "):
            if current_heading is not None:
                sections.append((current_heading, current_block))
            current_heading = line[3:].strip()
            current_block = []
            continue
        if current_heading is not None:
            current_block.append(line)

    if current_heading is not None:
        sections.append((current_heading, current_block))

    parsed: list[SlideSection] = []
    for heading, block in sections:
        title = heading
        body_lines: list[str] = []

        for raw in block:
            line = raw.strip()
            if not line or line == "---":
                continue

            if line.lower().startswith("**title:**"):
                title = _clean_line(line.split(":", 1)[1])
                continue

            if line.startswith("### "):
                body_lines.append(_clean_line(line[4:]))
                continue

            if line.startswith("- "):
                body_lines.append(_clean_line(line[2:]))
                continue

            if re.match(r"^\d+\.\s+", line):
                body_lines.append(_clean_line(re.sub(r"^\d+\.\s+", "", line)))
                continue

            if line.startswith("|"):
                body_lines.append(_clean_line(line))
                continue

            body_lines.append(_clean_line(line))

        # Remove markdown table delimiter rows and collapse duplicates.
        filtered: list[str] = []
        seen = set()
        for line in body_lines:
            if re.match(r"^\|?\s*[-: ]+\|", line):
                continue
            if line and line not in seen:
                filtered.append(line)
                seen.add(line)

        parsed.append(SlideSection(heading=heading, title=title, body_lines=filtered))

    return parsed


def build_presentation(sections: list[SlideSection], output_path: Path) -> None:
    prs = Presentation()

    # Title slide based on first content slide.
    first = sections[0] if sections else SlideSection("", "Aura-Discovery", ["DDI Prediction"]) 
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = first.title
    subtitle = slide.placeholders[1]
    subtitle.text = "Generated from project_ppt_full_deck.md"

    for section in sections[1:]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section.title

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        lines = section.body_lines[:14]
        if not lines:
            lines = [section.heading]

        for idx, line in enumerate(lines):
            para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            para.text = line
            para.level = 0
            para.font.size = Pt(18)

    prs.save(output_path)


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    source_path = repo_root / "reports" / "project_ppt_full_deck.md"
    output_path = repo_root / "reports" / "project_ppt_full_deck.pptx"

    markdown_text = source_path.read_text(encoding="utf-8")
    sections = parse_markdown_sections(markdown_text)
    build_presentation(sections, output_path)
    print(f"Created PPTX: {output_path}")
    print(f"Slides created: {len(sections)}")


if __name__ == "__main__":
    main()
